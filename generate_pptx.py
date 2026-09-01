import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_COLOR = RGBColor(11, 19, 43)        # Deep Navy / Dark Slate
    CARD_BG = RGBColor(22, 33, 62)         # Slate Blue Card
    CARD_BORDER = RGBColor(38, 70, 112)    # Border
    CYAN = RGBColor(6, 182, 212)           # Sovereign Cyan
    CYAN_LIGHT = RGBColor(103, 232, 249)   # Light Cyan
    EMERALD = RGBColor(16, 185, 129)       # Emerald Green
    WHITE = RGBColor(248, 250, 252)        # Bright White
    GRAY = RGBColor(148, 163, 184)         # Slate Gray
    AMBER = RGBColor(245, 158, 11)         # Amber Warning

    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="SOVEREIGN AI WORKBENCH • SIH26117"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Consolas"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    def add_card(slide, left, top, width, height, title=None, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = CYAN_LIGHT
        return card

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)
    
    h_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(3.0))
    tf1 = h_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON • PROBLEM STATEMENT ID: SIH26117"
    p1.font.name = "Consolas"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = CYAN
    p1.space_after = Pt(10)
    
    p2 = tf1.add_paragraph()
    p2.text = "SovereignAI Workbench"
    p2.font.name = "Arial"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "On-Premise, Air-Gapped Industrial AI Platform with Zero Cloud Exfiltration"
    p3.font.name = "Arial"
    p3.font.size = Pt(18)
    p3.font.color.rgb = GRAY
    p3.space_after = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = '“Own the Data. Own the Models. Own the Execution.”'
    p4.font.name = "Consolas"
    p4.font.size = Pt(16)
    p4.font.italic = True
    p4.font.color.rgb = EMERALD

    add_card(s1, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2))
    foot_box = s1.shapes.add_textbox(Inches(1.2), Inches(5.65), Inches(10.9), Inches(0.9))
    tff = foot_box.text_frame
    pf1 = tff.paragraphs[0]
    pf1.text = "DOMAIN: Critical Infrastructure, Defence, PSUs & Industrial Data Sovereignty"
    pf1.font.name = "Arial"
    pf1.font.size = Pt(13)
    pf1.font.bold = True
    pf1.font.color.rgb = WHITE
    
    pf2 = tff.add_paragraph()
    pf2.text = "TECH: Local Open-Weights LLMs (Llama 3.1 / DeepSeek R1 / Qwen 2.5) • TriForge Router • Local RAG • Python Sandbox • SHA-256 Audit"
    pf2.font.name = "Consolas"
    pf2.font.size = Pt(11)
    pf2.font.color.rgb = CYAN_LIGHT

    # ==========================================
    # SLIDE 2: Problem Statement
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "Problem Statement: Industrial AI & The Cloud Privacy Crisis")
    
    col_w = Inches(3.6)
    add_card(s2, Inches(0.8), Inches(1.5), col_w, Inches(5.3), "1. Cloud Data Exfiltration Risk", AMBER)
    tb = s2.shapes.add_textbox(Inches(0.95), Inches(2.1), col_w - Inches(0.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets1 = [
        "Modern industrial enterprises cannot use OpenAI / Claude / Cloud APIs due to strict data secrecy.",
        "Uploading plant blueprints, SCADA telemetry, and defence SOPs breaches air-gap isolation.",
        "Risk of IP theft, foreign intelligence interception, and national security compromise."
    ]
    for b in bullets1:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)

    add_card(s2, Inches(4.8), Inches(1.5), col_w, Inches(5.3), "2. Regulatory Mandates", CYAN)
    tb = s2.shapes.add_textbox(Inches(4.95), Inches(2.1), col_w - Inches(0.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets2 = [
        "DPDP Act (India), ISO 27001, and Defence Cybersecurity Directives mandate zero data egress.",
        "PSUs, nuclear plants, and refineries must operate completely offline in isolated air-gaps.",
        "No network socket to public DNS or external IP addresses is permitted."
    ]
    for b in bullets2:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)

    add_card(s2, Inches(8.8), Inches(1.5), col_w, Inches(5.3), "3. Autonomous Risk & Hallucination", EMERALD)
    tb = s2.shapes.add_textbox(Inches(8.95), Inches(2.1), col_w - Inches(0.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets3 = [
        "Unverified AI outputs can cause catastrophic physical machinery damage if hallucinated.",
        "Lack of sandboxed script execution and human-in-the-loop approval gates in current LLM tools.",
        "Zero traceability or tamper-proof audit trails for safety regulatory inspections."
    ]
    for b in bullets3:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: Existing Gap
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "Existing Gaps: Why Current Solutions Fail in High-Security Industries")

    rows, cols = 5, 4
    left, top, width, height = Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0)
    table_shape = s3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.9)
    table.columns[2].width = Inches(3.3)
    table.columns[3].width = Inches(3.3)

    headers = ["Approach", "Examples", "Critical Limitation", "SovereignAI Advantage"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = CYAN_LIGHT

    data = [
        ("Cloud AI APIs", "ChatGPT, Claude 3.5, Azure OpenAI", "100% cloud network egress; leaks proprietary blueprints & defence data.", "100% Air-Gapped on-premise local inference with 0 cloud dependencies."),
        ("Generic Local WebUIs", "Ollama CLI, OpenWebUI, LM Studio", "Just a raw chat prompt. No workflow graph, no approvals, no report deliverables.", "Turnkey agentic workflow: RAG + Python Sandbox + HITL + DOCX Reports."),
        ("Basic RAG Demos", "LangChain / LlamaIndex Scripts", "Prone to hallucinations; no citation verification; no tamper-proof audit trail.", "Multi-document OCR, page citation tracking, and SHA-256 verifiable audit ledger."),
        ("Legacy SCADA / ERP", "SAP, SCADA, Custom Portals", "Static rule-based systems; cannot reason over unstructured engineering PDFs.", "Multimodal Vision & Reasoning models analyze diagrams, P&IDs, and logs.")
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE if col_idx == 3 else (CYAN_LIGHT if col_idx == 0 else GRAY)
            if col_idx == 3:
                p.font.bold = True

    # ==========================================
    # SLIDE 4: Our Solution
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "Our Solution: SovereignAI Workbench Architecture")

    card_w = Inches(5.6)
    card_h = Inches(2.5)
    
    add_card(s4, Inches(0.8), Inches(1.5), card_w, card_h, "🛡️ Air-Gapped Local Inference Engine", CYAN)
    tb = s4.shapes.add_textbox(Inches(0.95), Inches(2.0), card_w - Inches(0.3), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for text in ["Zero cloud telemetry or external network calls.", "Runs open-weight models locally via Ollama / vLLM runtime.", "Hardware-aware execution tuned for enterprise workstation GPUs."]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    add_card(s4, Inches(6.8), Inches(1.5), card_w, card_h, "⚙️ TriForge Smart Model Router", EMERALD)
    tb = s4.shapes.add_textbox(Inches(6.95), Inches(2.0), card_w - Inches(0.3), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for text in ["Intelligent intent & task classification matrix.", "Routes to specialized models: DeepSeek R1 (Reasoning), Qwen 2.5 (Coding), Qwen 2 VL (Vision).", "Evaluates policy and security risk before dispatching."]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    add_card(s4, Inches(0.8), Inches(4.3), card_w, card_h, "🔒 Sandboxed Execution & HITL Approvals", AMBER)
    tb = s4.shapes.add_textbox(Inches(0.95), Inches(4.8), card_w - Inches(0.3), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for text in ["Isolated Python subprocess sandbox for numerical calculations.", "Human-in-the-Loop (HITL) gate halts dangerous system modifications.", "Operator approves script execution with single-click sign-off."]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    add_card(s4, Inches(6.8), Inches(4.3), card_w, card_h, "📑 Verifiable Deliverables & SHA-256 Audit", CYAN_LIGHT)
    tb = s4.shapes.add_textbox(Inches(6.95), Inches(4.8), card_w - Inches(0.3), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for text in ["Generates verified Word (.docx) compliance & approval reports.", "Citation verification engine validates claims against local SOPs.", "Tamper-evident audit ledger cryptographically logs every action."]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 5: Key Features
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "Key Features: 8 Architectural Pillars of SovereignAI")

    p_w, p_h = Inches(2.7), Inches(2.5)
    pillars = [
        ("1. TriForge Router", "Deterministic routing matrix based on classification & GPU VRAM.", CYAN),
        ("2. Local Vector RAG", "Confidential embeddings for PDFs, DOCX, XLSX with page citations.", CYAN),
        ("3. Multimodal Vision", "OCR & diagram analysis for P&ID schematics and scanned sheets.", CYAN),
        ("4. Graph Orchestrator", "Multi-step agent loop: Planning, Retrieval, Execution, Synthesis.", CYAN),
        ("5. Python Sandbox", "Isolated execution zone for pressure, maintenance, & math metrics.", EMERALD),
        ("6. HITL Approval Gate", "Operator sign-off inbox for high-risk actions before execution.", EMERALD),
        ("7. DOCX Deliverables", "Automated creation of signed, verified executive analysis notes.", EMERALD),
        ("8. SHA-256 Audit Trail", "Immutable ledger recording model calls, citations, and tool actions.", EMERALD),
    ]

    for idx, (title, desc, col) in enumerate(pillars):
        r = idx // 4
        c = idx % 4
        x = Inches(0.8) + c * Inches(2.95)
        y = Inches(1.5) + r * Inches(2.8)
        add_card(s5, x, y, p_w, p_h, title, col)
        tb = s5.shapes.add_textbox(x + Inches(0.15), y + Inches(0.65), p_w - Inches(0.3), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 6: TriForge Smart Model Router
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "TriForge: Policy-Aware Intelligent Model Routing Matrix")

    add_card(s6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Why Specialized Routing?", CYAN)
    tb = s6.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    t_points = [
        ("No Monolithic Compromise: ", "Generalist models are inefficient on local workstation GPUs. TriForge pairs specific tasks with specialized domain models."),
        ("Multi-Dimensional Inspection: ", "Evaluates task prompt keywords, modality (Text vs Vision), confidentiality level (Internal vs Restricted), and available GPU VRAM."),
        ("Latency & Cost Optimization: ", "Sub-second deterministic routing (<450ms) ensures maximum throughput without cloud API latency."),
        ("Safety Policy Check: ", "Flags sensitive operations (code execution, file modifications) before allocating GPU resources.")
    ]
    for h, b in t_points:
        p = tf.add_paragraph()
        p.text = "• " + h + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)

    add_card(s6, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "TriForge Routing Matrix", EMERALD)
    tb = s6.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    routes = [
        ("CODING & SCRIPTS", "Qwen 2.5 Coder 7B", "C++, Python, algorithm synthesis, debugging, script automation."),
        ("REASONING & SOPs", "DeepSeek R1 / Llama 3.1 8B", "Industrial compliance, deviation analysis, safety manual reasoning."),
        ("VISION & P&IDs", "Qwen 2 VL 7B", "Scanned drawings, equipment tags, engineering schematics OCR."),
        ("SPREADSHEETS & METRICS", "DeepSeek R1 + Sandbox", "Excel parsing, predictive maintenance score, pressure delta math.")
    ]
    for cat, model, desc in routes:
        p1 = tf.add_paragraph()
        p1.text = f"▶ {cat} → [{model}]"
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = CYAN_LIGHT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = GRAY
        p2.space_after = Pt(6)

    # ==========================================
    # SLIDE 7: System Architecture
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "System Architecture: End-to-End Air-Gapped Stack")

    layers = [
        ("1. Presentation Layer (React 18 / Vite)", "Sovereign Command Center UI • TriForge Routing Panel • HITL Approval Inbox • Live Audit Log Explorer", CYAN),
        ("2. Security & Policy Gateway (FastAPI)", "Network Sentinel Firewall (Blocks Egress) • Policy Engine • Zero Cloud Guardrails", AMBER),
        ("3. Industrial Orchestration Graph", "Task Classifier • Local Vector RAG Index • Verification Engine • DOCX Report Generator", CYAN_LIGHT),
        ("4. Execution & Sandbox Layer", "Isolated Python Calculation Sandbox • File Ingestion (PyMuPDF / docx / openpyxl)", EMERALD),
        ("5. On-Premise AI Model Runtime (Ollama)", "Llama 3.1 8B Instruct • DeepSeek R1 8B • Qwen 2.5 Coder 7B • Qwen 2 VL 7B", WHITE)
    ]

    for idx, (title, desc, col) in enumerate(layers):
        y = Inches(1.5) + idx * Inches(1.05)
        add_card(s7, Inches(0.8), y, Inches(11.7), Inches(0.95), None, col)
        tb = s7.shapes.add_textbox(Inches(1.0), y + Inches(0.08), Inches(11.3), Inches(0.75))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 8: How It Works
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "How It Works: 6-Step Sovereign Task Execution Lifecycle")

    steps = [
        ("Step 1: Input", "Operator submits industrial query or uploads inspection PDF/SOP."),
        ("Step 2: Policy & Route", "TriForge classifies intent, checks security policy, and selects local model."),
        ("Step 3: RAG Retrieval", "Confidential vector store retrieves top citations with exact page numbers."),
        ("Step 4: Real LLM Inference", "Local open-weights model synthesizes analysis & python script."),
        ("Step 5: HITL Approval", "Operator reviews calculation code in Approval Inbox and signs off."),
        ("Step 6: Verified Deliverable", "Sandbox executes code, verifies citations, and produces signed DOCX.")
    ]

    for idx, (s_title, s_desc) in enumerate(steps):
        r = idx // 3
        c = idx % 3
        x = Inches(0.8) + c * Inches(3.95)
        y = Inches(1.6) + r * Inches(2.7)
        add_card(s8, x, y, Inches(3.7), Inches(2.4), s_title, CYAN)
        tb = s8.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(3.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 9: Multimodal & Local RAG
    # ==========================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "Multimodal Document Intelligence & Confidential RAG")

    add_card(s9, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Zero-Cloud Ingestion Engine", CYAN)
    tb = s9.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    pts = [
        "Universal Format Support: Native local extraction for PDF, DOCX, XLSX, TXT, and scanned image records.",
        "High-Speed OCR: PyMuPDF + offline visual parser extracts text from low-quality scans and engineering stamps.",
        "Smart Dynamic Chunking: Preserves table hierarchies, section headings, and page boundaries.",
        "Prioritized Document Indexing: Newly uploaded inspection reports are prioritized over baseline SOP manuals."
    ]
    for p_text in pts:
        p = tf.add_paragraph()
        p.text = "• " + p_text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)

    add_card(s9, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Confidential RAG & Citation Guarantee", EMERALD)
    tb = s9.shapes.add_textbox(Inches(6.95), Inches(2.1), Inches(5.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    pts2 = [
        "100% On-Premise Vector Store: Embeddings reside entirely in workstation memory/disk without cloud vector DBs.",
        "Deterministic Grounding: Every answer includes source filename, chunk ID, and page references.",
        "Anti-Hallucination Guardrails: Claims are cross-referenced against SOP citations before deliverable generation.",
        "Session Isolation & Reset: Instant single-click workspace reset clears active document indices for new tasks."
    ]
    for p_text in pts2:
        p = tf.add_paragraph()
        p.text = "• " + p_text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 10: Security & Air-Gap Sovereignty
    # ==========================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, "Security & Sovereignty: Guaranteed Zero Cloud Egress")

    sec_cards = [
        ("🔒 Network Sentinel", "Monitors outbound network activity. Strictly blocks egress to external AI domains (OpenAI, Anthropic, HuggingFace). EXTERNAL_AI_CALLS = 0.", AMBER),
        ("🛡️ Deterministic Policy Engine", "Enforces access rules based on classification (CONFIDENTIAL, RESTRICTED). Requires operator authorization for tool use.", CYAN),
        ("📦 Isolated Python Sandbox", "Safe calculation environment with execution timeouts, restricted OS syscalls, and memory bounds.", EMERALD),
        ("📜 SHA-256 Audit Ledger", "Tamper-evident JSONL audit ledger records timestamp, user action, model used, prompt, and output hash.", CYAN_LIGHT)
    ]

    for idx, (title, desc, col) in enumerate(sec_cards):
        r = idx // 2
        c = idx % 2
        x = Inches(0.8) + c * Inches(5.95)
        y = Inches(1.6) + r * Inches(2.65)
        add_card(s10, x, y, Inches(5.7), Inches(2.4), title, col)
        tb = s10.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 11: Flagship Workflow
    # ==========================================
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s11)
    add_header(s11, "Flagship Workflow: Industrial Inspection to Verified DOCX Report")

    add_card(s11, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "End-to-End Pressure Safety Compliance Workflow", CYAN)
    tb = s11.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    wf_steps = [
        ("1. Ingestion: ", "Operator uploads 'TriForge_OnePager.pdf' / Inspection log into Sovereign AI Workbench."),
        ("2. Query: ", "'Analyze overpressure deviation against SOP-17, calculate delta %, and draft sign-off note.'"),
        ("3. Routing & RAG: ", "TriForge routes to DeepSeek R1 & retrieves SOP-17 Sec 4.2 (Ceiling: 120 PSI, Critical: 135 PSI)."),
        ("4. HITL Approval Gate: ", "Policy Engine flags Python execution as HIGH RISK. Dispatches ticket to Approval Inbox."),
        ("5. Operator Sign-Off: ", "Industrial Operator reviews formula (P_delta = (142.8 - 120) / 120) and clicks 'APPROVE'."),
        ("6. DOCX Generation: ", "Generates official 'Approval_Note.docx' with Executive Summary, Sandbox Output & Citations.")
    ]

    for title, desc in wf_steps:
        p = tf.add_paragraph()
        p.text = title + desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE if "4." in title or "6." in title else GRAY
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 12: Other Use Cases
    # ==========================================
    s12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s12)
    add_header(s12, "Broad-Spectrum Industrial & Enterprise Use Cases")

    uc = [
        ("💻 Software & C++ / Python", "Generates full C++ array code, sorting algorithms, PLC automation scripts, and firmware debug analysis.", CYAN),
        ("📊 Predictive Maintenance", "Parses equipment telemetry spreadsheets (.xlsx) and calculates MTBF and degradation scores.", EMERALD),
        ("📐 P&ID & Engineering Vision", "Extracts instrumentation tags, valve specs, and flow schematics from scanned blueprints.", CYAN_LIGHT),
        ("🛡️ Defence & PSU Compliance", "Instant search, summarization, and audit cross-referencing across confidential manuals.", AMBER)
    ]

    for idx, (title, desc, col) in enumerate(uc):
        r = idx // 2
        c = idx % 2
        x = Inches(0.8) + c * Inches(5.95)
        y = Inches(1.6) + r * Inches(2.65)
        add_card(s12, x, y, Inches(5.7), Inches(2.4), title, col)
        tb = s12.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 13: Technology Stack
    # ==========================================
    s13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s13)
    add_header(s13, "Technology Stack: Built on Modern, Sovereign Open Source")

    stack_cards = [
        ("Frontend UI", "React 18 • Vite • Tailwind CSS • Lucide Icons • Responsive Glassmorphism Dashboard", CYAN),
        ("Backend Services", "FastAPI (Python 3.11) • Pydantic v2 • Uvicorn • HTTPX Async Client", CYAN_LIGHT),
        ("Local AI Engine", "Ollama Local Runtime • Llama 3.1 8B • DeepSeek R1 • Qwen 2.5 Coder • Qwen 2 VL", EMERALD),
        ("Ingestion & Files", "PyMuPDF (fitz) • python-docx • openpyxl • NumPy • Subprocess Sandbox", AMBER)
    ]

    for idx, (title, desc, col) in enumerate(stack_cards):
        r = idx // 2
        c = idx % 2
        x = Inches(0.8) + c * Inches(5.95)
        y = Inches(1.6) + r * Inches(2.65)
        add_card(s13, x, y, Inches(5.7), Inches(2.4), title, col)
        tb = s13.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 14: Innovation & USPs
    # ==========================================
    s14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s14)
    add_header(s14, "Key Innovations & Unique Selling Propositions (USPs)")

    add_card(s14, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "What Sets SovereignAI Workbench Apart?", CYAN)
    tb = s14.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    usps = [
        ("1. Real Local Inference (No Mocks): ", "Directly integrates with local Ollama open-weights models. Explicit offline alerting with zero fake fallbacks."),
        ("2. Policy-Aware TriForge Routing: ", "Deterministically selects the optimal model based on data classification and task domain."),
        ("3. Mandatory HITL Safety Gate: ", "Prevents unmonitored code execution and system modifications in critical infrastructure."),
        ("4. Turnkey Verified Deliverables: ", "Automatically synthesizes signed Word (.docx) compliance reports with verified citations."),
        ("5. Cryptographic SHA-256 Audit Trail: ", "Built specifically for high-assurance defence and PSU compliance verification.")
    ]

    for title, desc in usps:
        p = tf.add_paragraph()
        p.text = title + desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE if "1." in title or "2." in title else GRAY
        p.space_after = Pt(12)

    # ==========================================
    # SLIDE 15: Demo & UI Walkthrough
    # ==========================================
    s15 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s15)
    add_header(s15, "User Interface & Operator Experience Walkthrough")

    ui_modules = [
        ("1. Command Center Dashboard", "Single-pane workbench featuring live Network Sentinel status, VRAM telemetry, and task input launcher.", CYAN),
        ("2. TriForge Model Router Panel", "Real-time visualization of model selection, risk rating, policy decision, and estimated latency.", CYAN_LIGHT),
        ("3. Agent Activity & State Trace", "Multi-step orchestration graph plan showing real-time execution steps and Markdown/code formatting.", EMERALD),
        ("4. Operator Approval Inbox", "Human-in-the-Loop gate showing proposed Python script, risk level, and Approve/Reject controls.", AMBER)
    ]

    for idx, (title, desc, col) in enumerate(ui_modules):
        r = idx // 2
        c = idx % 2
        x = Inches(0.8) + c * Inches(5.95)
        y = Inches(1.6) + r * Inches(2.65)
        add_card(s15, x, y, Inches(5.7), Inches(2.4), title, col)
        tb = s15.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 16: Results & Validation
    # ==========================================
    s16 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s16)
    add_header(s16, "Results, Validation & Performance Benchmarks")

    add_card(s16, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), "Empirical Test Results & System Validation", EMERALD)
    tb = s16.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    results = [
        ("✅ 42+ Automated Test Cases Passed: ", "Verified vector store ingestion, PyMuPDF OCR, RAG ranking, and FastAPI endpoints."),
        ("✅ Real Local Inference Verified: ", "Tested live with Llama 3.1 8B on local Ollama server with zero external network socket calls."),
        ("✅ C++ / Python Code Generation: ", "Generated complete, runnable C++ array programs with O(N) traversal and sorting."),
        ("✅ Accurate Document Summarization: ", "Extracted structured findings from 'TriForge_OnePager.pdf' citing exact parameters."),
        ("⚡ High-Speed Local Latency: ", "Sub-second model routing (<450ms) and fast local token generation on standard hardware.")
    ]

    for title, desc in results:
        p = tf.add_paragraph()
        p.text = title + desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE if "✅" in title else GRAY
        p.space_after = Pt(12)

    # ==========================================
    # SLIDE 17: Scalability & Future Scope
    # ==========================================
    s17 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s17)
    add_header(s17, "Scalability & Future Roadmap")

    future_cards = [
        ("🗄️ Enterprise Vector Store", "Integration with on-premise PostgreSQL + pgvector for multi-million document repositories.", CYAN),
        ("🎯 Domain LoRA Fine-Tuning", "Lightweight fine-tuning adapters for defence terminology and PSU technical standards.", CYAN_LIGHT),
        ("🏭 Edge Hardware Deployment", "Containerized deployment on ruggedized industrial edge devices (NVIDIA Jetson AGX).", EMERALD),
        ("🤖 Multi-Agent Debate", "Autonomous multi-agent consensus loops for mission-critical industrial decision validation.", AMBER)
    ]

    for idx, (title, desc, col) in enumerate(future_cards):
        r = idx // 2
        c = idx % 2
        x = Inches(0.8) + c * Inches(5.95)
        y = Inches(1.6) + r * Inches(2.65)
        add_card(s17, x, y, Inches(5.7), Inches(2.4), title, col)
        tb = s17.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 18: Impact & Value Proposition
    # ==========================================
    s18 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s18)
    add_header(s18, "Strategic Impact & Business Value Proposition")

    impacts = [
        ("🏛️ Defence & National Security", "Guarantees complete technological sovereignty with zero risk of cloud foreign surveillance.", CYAN),
        ("⏱️ 90% Time Reduction", "Reduces complex engineering SOP lookup and incident analysis time from 3 hours to under 30 seconds.", EMERALD),
        ("💰 Zero Cloud API OpEx", "Eliminates recurrent per-token subscription costs by utilizing existing on-premise GPU workstations.", CYAN_LIGHT),
        ("🛡️ Zero Compliance Liability", "100% compliant with DPDP Act, ISO 27001, and critical infrastructure air-gap mandates.", AMBER)
    ]

    for idx, (title, desc, col) in enumerate(impacts):
        r = idx // 2
        c = idx % 2
        x = Inches(0.8) + c * Inches(5.95)
        y = Inches(1.6) + r * Inches(2.65)
        add_card(s18, x, y, Inches(5.7), Inches(2.4), title, col)
        tb = s18.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = GRAY

    # ==========================================
    # SLIDE 19: Conclusion
    # ==========================================
    s19 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s19)
    
    h_box = s19.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf19 = h_box.text_frame
    tf19.word_wrap = True
    
    p1 = tf19.paragraphs[0]
    p1.text = "CONCLUSION & SUMMARY"
    p1.font.name = "Consolas"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = CYAN
    p1.space_after = Pt(10)
    
    p2 = tf19.add_paragraph()
    p2.text = "SovereignAI Workbench"
    p2.font.name = "Arial"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(15)

    pillars_text = [
        "1. OWN THE DATA — Zero cloud exfiltration; files never leave local storage.",
        "2. OWN THE MODELS — Full local execution over open-weights models (Ollama).",
        "3. OWN THE EXECUTION — Sandboxed calculations, HITL approvals, and signed DOCX deliverables."
    ]
    for p_t in pillars_text:
        p = tf19.add_paragraph()
        p.text = p_t
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN_LIGHT
        p.space_after = Pt(8)

    p_end = tf19.add_paragraph()
    p_end.text = "The Future of Critical Enterprise Industrial AI is 100% Sovereign."
    p_end.font.name = "Consolas"
    p_end.font.size = Pt(16)
    p_end.font.italic = True
    p_end.font.color.rgb = EMERALD

    # ==========================================
    # SLIDE 20: Thank You / Q&A
    # ==========================================
    s20 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s20)
    
    add_card(s20, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5), None, CYAN)
    tb = s20.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.7), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "SovereignAI Workbench • Smart India Hackathon (SIH26117)"
    p2.font.name = "Consolas"
    p2.font.size = Pt(16)
    p2.font.color.rgb = CYAN_LIGHT
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "GitHub Repository: https://github.com/Sarthak752008/sovereign-ai-workbench"
    p3.font.name = "Consolas"
    p3.font.size = Pt(13)
    p3.font.color.rgb = EMERALD
    p3.space_after = Pt(10)

    p4 = tf.add_paragraph()
    p4.text = "Ready for Live System Demonstration & Jury Q&A"
    p4.font.name = "Arial"
    p4.font.size = Pt(18)
    p4.font.bold = True
    p4.font.color.rgb = WHITE

    output_path = os.path.join(r"c:\Users\sarth\OneDrive\Desktop\sovereign-ai-workbench", "SovereignAI_Workbench_SIH26117_Presentation.pptx")
    prs.save(output_path)
    print(f"SUCCESS: Presentation created at {output_path}")

if __name__ == "__main__":
    create_presentation()
